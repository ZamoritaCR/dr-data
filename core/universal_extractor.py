"""
Universal File Extractor — Sprint 2A.

Extracts structured content from any supported file type into a
unified dict format that feeds the dashboard intelligence layer.

Supported: PDF, DOCX, DOC, PPTX, PPT, JSON, TXT, CSV, XLSX, XLS, TWB, TWBX.
"""

from __future__ import annotations

import csv
import io
import json
import os
import re
from pathlib import Path
from typing import Any


# ---------------------------------------------------------------------------
#  Helpers
# ---------------------------------------------------------------------------

_NUM_RE = re.compile(
    r"(?<![A-Za-z/\\])"            # not preceded by letter or path sep
    r"[$€£¥]?\s*"                  # optional currency symbol
    r"-?"                          # optional negative
    r"\d{1,3}(?:[,_]\d{3})*"      # integer part with optional grouping
    r"(?:\.\d+)?"                  # optional decimal
    r"%?"                          # optional percent
    r"(?!\.\d)"                    # not followed by more decimals (version strings)
)

_UNIT_RE = re.compile(
    r"(?:USD|EUR|GBP|JPY|CAD|AUD|BRL|MXN|CRC|COP|"
    r"kg|lbs?|mi|km|hrs?|mins?|secs?|days?|months?|years?|"
    r"units?|items?|rows?|records?|users?|visits?|"
    r"MB|GB|TB|KB|bps|Mbps|Gbps|%)",
    re.IGNORECASE,
)


def _parse_number(raw: str) -> float | None:
    """Try to parse a numeric string, stripping currency and grouping."""
    s = raw.strip().lstrip("$€£¥").replace(",", "").replace("_", "").rstrip("%")
    try:
        return float(s)
    except (ValueError, OverflowError):
        return None


def _extract_numbers_from_text(text: str, max_items: int = 200) -> list[dict]:
    """Pull numbers with surrounding context from raw text."""
    results: list[dict] = []
    for m in _NUM_RE.finditer(text):
        val = _parse_number(m.group())
        if val is None:
            continue
        start = max(0, m.start() - 60)
        end = min(len(text), m.end() + 60)
        context = text[start:end].replace("\n", " ").strip()
        unit_m = _UNIT_RE.search(context)
        results.append({
            "value": val,
            "context": context,
            "unit": unit_m.group() if unit_m else "",
        })
        if len(results) >= max_items:
            break
    return results


def _infer_field_type(values: list) -> str:
    """Guess a column type from sample values."""
    non_none = [v for v in values if v is not None and str(v).strip() != ""]
    if not non_none:
        return "unknown"
    nums = 0
    dates = 0
    for v in non_none[:20]:
        s = str(v).strip()
        if _parse_number(s) is not None:
            nums += 1
        if re.match(r"\d{4}[-/]\d{1,2}([-/]\d{1,2})?", s):
            dates += 1
    ratio = len(non_none[:20])
    if dates / max(ratio, 1) > 0.6:
        return "date"
    if nums / max(ratio, 1) > 0.6:
        return "numeric"
    avg_len = sum(len(str(v)) for v in non_none[:20]) / max(len(non_none[:20]), 1)
    if avg_len > 80:
        return "text_long"
    return "text"


def _infer_schema(records: list[dict], max_sample: int = 5) -> dict:
    """Infer schema from a list of flat dicts."""
    if not records:
        return {"fields": []}
    fields: list[dict] = []
    keys = list(records[0].keys()) if isinstance(records[0], dict) else []
    for k in keys:
        vals = [r.get(k) for r in records[:50] if isinstance(r, dict)]
        sample = [v for v in vals[:max_sample] if v is not None]
        fields.append({
            "name": k,
            "type": _infer_field_type(vals),
            "sample_values": [str(v) for v in sample],
        })
    return {"fields": fields}


def _schema_from_table(headers: list[str], rows: list[list]) -> dict:
    """Build schema hint from a table with headers and rows."""
    if not headers:
        return {"fields": []}
    fields: list[dict] = []
    for i, h in enumerate(headers):
        vals = [r[i] for r in rows[:50] if i < len(r)]
        sample = [str(v) for v in vals[:5] if v is not None]
        fields.append({
            "name": str(h),
            "type": _infer_field_type(vals),
            "sample_values": sample,
        })
    return {"fields": fields}


def _extract_headers_from_text(text: str) -> list[str]:
    """Heuristic: lines that look like section headers."""
    headers: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped or len(stripped) > 120:
            continue
        # ALL CAPS lines, markdown headings, numbered headings
        if (
            stripped.isupper() and len(stripped) > 3
            or stripped.startswith("#")
            or re.match(r"^\d+[\.\)]\s+[A-Z]", stripped)
            or (stripped.endswith(":") and len(stripped) < 80)
        ):
            clean = stripped.lstrip("#").strip().rstrip(":")
            if clean:
                headers.append(clean)
    return headers[:100]


def _safe(fn):
    """Decorator: catch all exceptions, return error dict instead of raising."""
    def wrapper(self, path: str) -> dict:
        try:
            return fn(self, path)
        except Exception as exc:
            return _error_result(path, str(exc))
    return wrapper


def _error_result(path: str, msg: str) -> dict:
    return {
        "file_type": Path(path).suffix.lower(),
        "file_name": Path(path).name,
        "raw_text": "",
        "tables": [],
        "numbers": [],
        "headers": [],
        "charts": [],
        "has_data": False,
        "schema_hint": {"fields": []},
        "extraction_method": "error",
        "error": msg,
    }


# ---------------------------------------------------------------------------
#  Main class
# ---------------------------------------------------------------------------

class UniversalExtractor:
    """Extracts structured content from any file type for dashboard generation."""

    SUPPORTED_TYPES = [
        ".pdf", ".docx", ".doc", ".pptx", ".ppt",
        ".json", ".txt", ".csv", ".xlsx", ".xls",
        ".twb", ".twbx",
    ]

    _DISPATCH: dict[str, str] = {
        ".pdf":  "_extract_pdf",
        ".docx": "_extract_docx",
        ".doc":  "_extract_docx",
        ".pptx": "_extract_pptx",
        ".ppt":  "_extract_pptx",
        ".json": "_extract_json",
        ".txt":  "_extract_txt",
        ".csv":  "_extract_csv_xlsx",
        ".xlsx": "_extract_csv_xlsx",
        ".xls":  "_extract_csv_xlsx",
        ".twb":  "_extract_twb",
        ".twbx": "_extract_twb",
    }

    # ------------------------------------------------------------------ #
    #  Public API                                                         #
    # ------------------------------------------------------------------ #

    def extract(self, file_path: str) -> dict:
        """Extract structured content from any supported file.

        Returns a unified dict with raw_text, tables, numbers, headers,
        charts, schema_hint, and metadata.  Never raises — errors are
        returned inside the dict.
        """
        path = str(file_path)
        if not os.path.isfile(path):
            return _error_result(path, f"File not found: {path}")

        ext = self.detect_file_type(path)
        method_name = self._DISPATCH.get(ext)
        if method_name is None:
            return _error_result(path, f"Unsupported file type: {ext}")

        method = getattr(self, method_name)
        return method(path)

    @classmethod
    def detect_file_type(cls, path: str) -> str:
        """Return the normalised file extension (e.g. '.pdf')."""
        return Path(path).suffix.lower()

    # ------------------------------------------------------------------ #
    #  PDF                                                                #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_pdf(self, path: str) -> dict:
        tables: list[dict] = []
        raw_parts: list[str] = []

        # --- pdfplumber for tables ---
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    raw_parts.append(text)
                    for tbl in (page.extract_tables() or []):
                        if not tbl or not tbl[0]:
                            continue
                        hdrs = [str(c or "") for c in tbl[0]]
                        rows = [[str(c or "") for c in r] for r in tbl[1:]]
                        tables.append({"headers": hdrs, "rows": rows})
        except Exception:
            pass

        # --- PyMuPDF fallback for text if pdfplumber gave nothing ---
        if not raw_parts or all(not t.strip() for t in raw_parts):
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(path)
                raw_parts = [page.get_text() for page in doc]
                doc.close()
            except Exception:
                pass

        raw_text = "\n".join(raw_parts)
        numbers = _extract_numbers_from_text(raw_text)
        headers = _extract_headers_from_text(raw_text)
        schema = (
            _schema_from_table(tables[0]["headers"], tables[0]["rows"])
            if tables
            else _infer_schema(
                [{"value": n["value"], "context": n["context"]} for n in numbers[:20]]
            )
        )

        return {
            "file_type": ".pdf",
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": numbers,
            "headers": headers,
            "charts": [],
            "has_data": bool(tables or numbers),
            "schema_hint": schema,
            "extraction_method": "pdfplumber+pymupdf",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  DOCX / DOC                                                        #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_docx(self, path: str) -> dict:
        from docx import Document

        doc = Document(path)
        raw_parts: list[str] = []
        headers: list[str] = []
        tables: list[dict] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            raw_parts.append(text)
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                headers.append(text)

        for tbl in doc.tables:
            rows_data: list[list[str]] = []
            for row in tbl.rows:
                rows_data.append([cell.text.strip() for cell in row.cells])
            if rows_data:
                hdrs = rows_data[0]
                tables.append({"headers": hdrs, "rows": rows_data[1:]})

        raw_text = "\n".join(raw_parts)
        numbers = _extract_numbers_from_text(raw_text)
        schema = (
            _schema_from_table(tables[0]["headers"], tables[0]["rows"])
            if tables
            else {"fields": []}
        )

        return {
            "file_type": Path(path).suffix.lower(),
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": numbers,
            "headers": headers,
            "charts": [],
            "has_data": bool(tables or numbers),
            "schema_hint": schema,
            "extraction_method": "python-docx",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  PPTX / PPT                                                        #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_pptx(self, path: str) -> dict:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401

        prs = Presentation(path)
        raw_parts: list[str] = []
        headers: list[str] = []
        tables: list[dict] = []
        charts: list[dict] = []

        for slide_idx, slide in enumerate(prs.slides):
            # Title
            if slide.shapes.title and slide.shapes.title.text.strip():
                headers.append(slide.shapes.title.text.strip())

            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            raw_parts.append(t)

                # Tables
                if shape.has_table:
                    tbl = shape.table
                    rows_data: list[list[str]] = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    if rows_data:
                        tables.append({
                            "headers": rows_data[0],
                            "rows": rows_data[1:],
                        })

                # Charts — extract series data
                if shape.has_chart:
                    chart = shape.chart
                    chart_info: dict[str, Any] = {
                        "title": "",
                        "type": str(chart.chart_type) if chart.chart_type else "unknown",
                        "series_data": [],
                    }
                    try:
                        if chart.chart_title and chart.chart_title.has_text_frame:
                            chart_info["title"] = chart.chart_title.text_frame.text
                    except Exception:
                        pass

                    for series in chart.series:
                        sd: dict[str, Any] = {"name": "", "values": []}
                        try:
                            sd["name"] = series.tx.strRef.strCache.pt[0].v if series.tx else ""
                        except Exception:
                            pass
                        try:
                            sd["values"] = [
                                float(pt.v) for pt in series.values
                                if pt.v is not None
                            ]
                        except Exception:
                            try:
                                sd["values"] = list(series.values)
                            except Exception:
                                pass
                        chart_info["series_data"].append(sd)

                    charts.append(chart_info)

        raw_text = "\n".join(raw_parts)
        numbers = _extract_numbers_from_text(raw_text)
        schema = (
            _schema_from_table(tables[0]["headers"], tables[0]["rows"])
            if tables
            else {"fields": []}
        )

        return {
            "file_type": Path(path).suffix.lower(),
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": numbers,
            "headers": headers,
            "charts": charts,
            "has_data": bool(tables or charts or numbers),
            "schema_hint": schema,
            "extraction_method": "python-pptx",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  JSON                                                              #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_json(self, path: str) -> dict:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)

        raw_text = json.dumps(data, indent=2, default=str)[:50_000]
        tables: list[dict] = []
        records: list[dict] = []

        # Array of objects → table
        if isinstance(data, list) and data and isinstance(data[0], dict):
            records = data
            hdrs = list(data[0].keys())
            rows = [[str(r.get(k, "")) for k in hdrs] for r in data[:500]]
            tables.append({"headers": hdrs, "rows": rows})

        # Single object with list values → try pivoting
        elif isinstance(data, dict):
            list_keys = [k for k, v in data.items() if isinstance(v, list)]
            if list_keys:
                for lk in list_keys:
                    arr = data[lk]
                    if arr and isinstance(arr[0], dict):
                        hdrs = list(arr[0].keys())
                        rows = [[str(r.get(k, "")) for k in hdrs] for r in arr[:500]]
                        tables.append({"headers": hdrs, "rows": rows})
                        records.extend(arr)
            else:
                # Flat object
                records = [data]
                hdrs = list(data.keys())
                rows = [[str(data.get(k, "")) for k in hdrs]]
                tables.append({"headers": hdrs, "rows": rows})

        numbers = _extract_numbers_from_text(raw_text)
        schema = _infer_schema(records[:50]) if records else {"fields": []}

        return {
            "file_type": ".json",
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": numbers,
            "headers": list(data.keys()) if isinstance(data, dict) else [],
            "charts": [],
            "has_data": bool(tables),
            "schema_hint": schema,
            "extraction_method": "json_stdlib",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  TXT                                                               #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_txt(self, path: str) -> dict:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw_text = f.read()

        # Try to detect CSV-like structure
        tables: list[dict] = []
        lines = raw_text.splitlines()
        if lines:
            sniffer = csv.Sniffer()
            try:
                dialect = sniffer.sniff(lines[0])
                reader = csv.reader(io.StringIO(raw_text), dialect)
                all_rows = list(reader)
                if all_rows and len(all_rows[0]) > 1:
                    tables.append({
                        "headers": all_rows[0],
                        "rows": all_rows[1:],
                    })
            except csv.Error:
                pass

        numbers = _extract_numbers_from_text(raw_text)
        headers = _extract_headers_from_text(raw_text)
        schema = (
            _schema_from_table(tables[0]["headers"], tables[0]["rows"])
            if tables
            else {"fields": []}
        )

        return {
            "file_type": ".txt",
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": numbers,
            "headers": headers,
            "charts": [],
            "has_data": bool(tables or numbers),
            "schema_hint": schema,
            "extraction_method": "text_stdlib",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  CSV / XLSX / XLS                                                  #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_csv_xlsx(self, path: str) -> dict:
        import pandas as pd

        ext = Path(path).suffix.lower()
        tables: list[dict] = []
        all_records: list[dict] = []

        if ext == ".csv":
            df = pd.read_csv(path, nrows=2000, encoding_errors="replace")
            hdrs = [str(c) for c in df.columns]
            rows = df.astype(str).values.tolist()
            tables.append({"headers": hdrs, "rows": rows})
            all_records = df.head(50).to_dict(orient="records")
        else:
            xls = pd.ExcelFile(path)
            for sheet in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet, nrows=2000)
                hdrs = [str(c) for c in df.columns]
                rows = df.astype(str).values.tolist()
                tables.append({"headers": hdrs, "rows": rows})
                all_records.extend(df.head(50).to_dict(orient="records"))

        raw_text = "\n".join(
            ", ".join(t["headers"]) + "\n" +
            "\n".join(", ".join(r) for r in t["rows"][:20])
            for t in tables
        )
        numbers = _extract_numbers_from_text(raw_text)
        schema = _infer_schema(all_records[:50]) if all_records else (
            _schema_from_table(tables[0]["headers"], tables[0]["rows"])
            if tables
            else {"fields": []}
        )

        return {
            "file_type": ext,
            "file_name": Path(path).name,
            "raw_text": raw_text[:50_000],
            "tables": tables,
            "numbers": numbers,
            "headers": [t["headers"] for t in tables] if tables else [],
            "charts": [],
            "has_data": bool(tables),
            "schema_hint": schema,
            "extraction_method": "pandas",
            "error": None,
        }

    # ------------------------------------------------------------------ #
    #  TWB / TWBX (delegate to existing parser)                          #
    # ------------------------------------------------------------------ #

    @_safe
    def _extract_twb(self, path: str) -> dict:
        # Try to use the existing enhanced parser first
        spec: dict = {}
        method = "xml_fallback"
        try:
            from core.enhanced_tableau_parser import parse_twb
            spec = parse_twb(path)
            method = "enhanced_tableau_parser"
        except ImportError:
            pass

        if not spec:
            # Minimal XML extraction
            import defusedxml.ElementTree as ET
            import zipfile, tempfile, shutil

            actual = path
            tmpdir = None
            if path.lower().endswith(".twbx"):
                tmpdir = tempfile.mkdtemp(prefix="ue_twbx_")
                with zipfile.ZipFile(path, "r") as z:
                    for name in z.namelist():
                        if name.endswith(".twb"):
                            z.extract(name, tmpdir)
                            actual = os.path.join(tmpdir, name)
                            break

            try:
                tree = ET.parse(actual)
                root = tree.getroot()
                spec = {
                    "type": "tableau",
                    "version": root.attrib.get("version", ""),
                    "datasources": [
                        ds.attrib.get("caption", ds.attrib.get("name", ""))
                        for ds in root.iter("datasource")
                    ],
                    "worksheets": [
                        ws.attrib.get("name", "")
                        for ws in root.iter("worksheet")
                    ],
                }
            finally:
                if tmpdir:
                    shutil.rmtree(tmpdir, ignore_errors=True)

        raw_text = json.dumps(spec, indent=2, default=str)[:50_000]
        worksheets = spec.get("worksheets", [])
        datasources = spec.get("datasources", [])
        calc_fields = spec.get("calculated_fields", [])

        # Build a pseudo-table of worksheet→datasource mapping
        tables: list[dict] = []
        if worksheets:
            tables.append({
                "headers": ["worksheet", "type"],
                "rows": [[str(w) if isinstance(w, str) else w.get("name", ""), "sheet"] for w in worksheets],
            })

        schema_fields = []
        if calc_fields:
            for cf in calc_fields[:20]:
                name = cf.get("name", cf) if isinstance(cf, dict) else str(cf)
                schema_fields.append({
                    "name": str(name),
                    "type": "calculated",
                    "sample_values": [],
                })

        return {
            "file_type": Path(path).suffix.lower(),
            "file_name": Path(path).name,
            "raw_text": raw_text,
            "tables": tables,
            "numbers": [],
            "headers": [str(w) if isinstance(w, str) else w.get("name", "") for w in worksheets],
            "charts": [],
            "has_data": bool(worksheets or datasources),
            "schema_hint": {"fields": schema_fields} if schema_fields else {"fields": []},
            "extraction_method": method,
            "error": None,
        }


# ---------------------------------------------------------------------------
#  Self-test
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import tempfile

    ext = UniversalExtractor()

    # --- Test 1: JSON array ---
    test_data = [
        {"date": "2026-01", "revenue": 125000, "region": "West"},
        {"date": "2026-02", "revenue": 138000, "region": "East"},
    ]
    with tempfile.NamedTemporaryFile(suffix=".json", mode="w", delete=False) as f:
        json.dump(test_data, f)
        fname = f.name
    result = ext.extract(fname)
    t1 = result["has_data"] and not result["error"]
    print(f"TEST 1 (JSON array):  {'PASS' if t1 else 'FAIL'}")
    print(f"  Schema hint: {result['schema_hint']}")
    os.unlink(fname)

    # --- Test 2: CSV ---
    csv_content = "name,value,date\nAlpha,100,2026-01\nBeta,200,2026-02\nGamma,300,2026-03\n"
    with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False) as f:
        f.write(csv_content)
        fname = f.name
    result = ext.extract(fname)
    t2 = result["has_data"] and len(result["tables"]) == 1 and not result["error"]
    print(f"TEST 2 (CSV):         {'PASS' if t2 else 'FAIL'}")
    print(f"  Tables: {len(result['tables'])}, Schema fields: {len(result['schema_hint']['fields'])}")
    os.unlink(fname)

    # --- Test 3: TXT with numbers ---
    txt_content = "Q1 Revenue: $1,250,000\nQ2 Revenue: $1,380,000\nGrowth: 10.4%\n"
    with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False) as f:
        f.write(txt_content)
        fname = f.name
    result = ext.extract(fname)
    t3 = len(result["numbers"]) >= 3 and not result["error"]
    print(f"TEST 3 (TXT numbers): {'PASS' if t3 else 'FAIL'}")
    print(f"  Numbers found: {len(result['numbers'])}")
    os.unlink(fname)

    # --- Test 4: JSON single object ---
    obj_data = {"company": "TAOP", "revenue": 500000, "employees": 15}
    with tempfile.NamedTemporaryFile(suffix=".json", mode="w", delete=False) as f:
        json.dump(obj_data, f)
        fname = f.name
    result = ext.extract(fname)
    t4 = result["has_data"] and not result["error"]
    print(f"TEST 4 (JSON object): {'PASS' if t4 else 'FAIL'}")
    os.unlink(fname)

    # --- Test 5: Missing file ---
    result = ext.extract("/tmp/does_not_exist_12345.pdf")
    t5 = result["error"] is not None and not result["has_data"]
    print(f"TEST 5 (missing file): {'PASS' if t5 else 'FAIL'}")

    # --- Test 6: detect_file_type ---
    t6 = (
        UniversalExtractor.detect_file_type("report.PDF") == ".pdf"
        and UniversalExtractor.detect_file_type("/tmp/data.xlsx") == ".xlsx"
    )
    print(f"TEST 6 (detect_type): {'PASS' if t6 else 'FAIL'}")

    # --- Summary ---
    all_pass = all([t1, t2, t3, t4, t5, t6])
    print(f"\n{'='*50}")
    print(f"EXTRACTOR TEST: {'PASS' if all_pass else 'FAIL'} ({sum([t1,t2,t3,t4,t5,t6])}/6)")
