"""
Export Engine -- generates PowerPoint, PDF, and Word from DataFrames.

Produces WU-branded, insight-driven reports (not generic stat dumps).
"""
import os
import tempfile
from datetime import datetime
from pathlib import Path

import pandas as pd
import numpy as np

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PDF
from reportlab.lib.pagesizes import letter
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# Word
from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH


# ------------------------------------------------------------------ #
#  WU Brand Colors                                                     #
# ------------------------------------------------------------------ #
WU_BLACK = "#0D0D0D"
WU_SURFACE = "#1A1A1A"
WU_ELEVATED = "#262626"
WU_BORDER = "#333333"
WU_YELLOW = "#FFE600"
WU_YELLOW2 = "#FFDE00"
WU_WHITE = "#FFFFFF"
WU_TEXT2 = "#B0B0B0"
WU_MUTED = "#808080"
WU_GREEN = "#238636"
WU_AMBER = "#d29922"
WU_RED = "#da3633"


# ------------------------------------------------------------------ #
#  Data Intelligence Helpers                                           #
# ------------------------------------------------------------------ #
def _classify_columns(df):
    """Classify columns into numeric, categorical, date, id-like."""
    numeric, categorical, date_cols, id_like = [], [], [], []
    for col in df.columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            if df[col].nunique() > len(df) * 0.9:
                id_like.append(col)
            else:
                numeric.append(col)
        elif pd.api.types.is_datetime64_any_dtype(df[col]):
            date_cols.append(col)
        else:
            nunique = df[col].nunique()
            if nunique <= 30:
                categorical.append(col)
            elif nunique > len(df) * 0.8:
                id_like.append(col)
    return {"numeric": numeric, "categorical": categorical,
            "date": date_cols, "id_like": id_like}


def _compute_insights(df):
    """Compute real, data-driven insights. Returns a dict of findings."""
    cols = _classify_columns(df)
    numeric = cols["numeric"]
    categorical = cols["categorical"]
    rows, total_cols = df.shape
    insights = {
        "rows": rows, "cols": total_cols,
        "numeric_count": len(numeric),
        "categorical_count": len(categorical),
        "null_pct": round(df.isnull().mean().mean() * 100, 1),
        "kpis": [],           # list of {label, total, avg, fmt_total, fmt_avg}
        "top_performers": [],  # list of {dimension, metric, top, top_val, top_pct, top3, top3_pct}
        "composition": [],     # list of {dimension, breakdown: [{name, val, pct}]}
        "findings": [],        # plain-text bullet points
        "quality_notes": [],   # data quality flags
    }

    # KPIs for top 4 numeric columns
    for col in numeric[:4]:
        s = pd.to_numeric(df[col], errors="coerce").dropna()
        if len(s) == 0:
            continue
        total = s.sum()
        avg = s.mean()
        insights["kpis"].append({
            "label": col.replace("_", " ").replace("-", " ").title(),
            "column": col,
            "total": total, "avg": avg,
            "min": s.min(), "max": s.max(), "median": s.median(),
            "fmt_total": _fmt(total), "fmt_avg": _fmt(avg),
        })

    # Top performers: first numeric by each categorical
    for cat in categorical[:3]:
        if not numeric:
            break
        metric = numeric[0]
        grouped = df.groupby(cat, dropna=True)[metric].sum().sort_values(ascending=False)
        if grouped.empty:
            continue
        total = grouped.sum()
        top_name = grouped.index[0]
        top_val = grouped.iloc[0]
        top_pct = round(top_val / total * 100, 1) if total else 0
        top3_names = list(grouped.head(3).index)
        top3_sum = grouped.head(3).sum()
        top3_pct = round(top3_sum / total * 100, 1) if total else 0
        insights["top_performers"].append({
            "dimension": cat.replace("_", " ").title(),
            "metric": metric.replace("_", " ").title(),
            "top": str(top_name), "top_val": _fmt(top_val), "top_pct": top_pct,
            "top3": [str(n) for n in top3_names], "top3_pct": top3_pct,
            "count": len(grouped),
        })

    # Composition breakdown for first low-card categorical
    low_card = [c for c in categorical if df[c].nunique() <= 8]
    for cat in low_card[:2]:
        if not numeric:
            break
        metric = numeric[0]
        grouped = df.groupby(cat, dropna=True)[metric].sum().sort_values(ascending=False)
        total = grouped.sum()
        breakdown = []
        for name, val in grouped.items():
            pct = round(val / total * 100, 1) if total else 0
            breakdown.append({"name": str(name), "val": _fmt(val), "pct": pct})
        insights["composition"].append({
            "dimension": cat.replace("_", " ").title(),
            "metric": metric.replace("_", " ").title(),
            "breakdown": breakdown,
        })

    # Generate plain-text findings
    if insights["kpis"]:
        k = insights["kpis"][0]
        insights["findings"].append(
            f"Total {k['label']}: {k['fmt_total']} "
            f"(average {k['fmt_avg']} across {rows:,} records)"
        )
    for tp in insights["top_performers"][:2]:
        insights["findings"].append(
            f"Top {tp['dimension']}: {tp['top']} accounts for "
            f"{tp['top_pct']}% of all {tp['metric']}"
        )
        if tp["top3_pct"] > 60:
            insights["findings"].append(
                f"High concentration: top 3 {tp['dimension']}s "
                f"({', '.join(tp['top3'])}) represent {tp['top3_pct']}% of total"
            )

    # Quality notes
    null_cols = [(c, round(df[c].isnull().mean() * 100, 1))
                 for c in df.columns if df[c].isnull().mean() > 0.05]
    if null_cols:
        insights["quality_notes"].append(
            f"{len(null_cols)} column(s) have >5% null values"
        )
    dup_count = df.duplicated().sum()
    if dup_count > 0:
        insights["quality_notes"].append(
            f"{dup_count:,} duplicate rows detected ({round(dup_count/rows*100, 1)}%)"
        )
    if not insights["quality_notes"]:
        insights["quality_notes"].append("No significant data quality issues detected")

    return insights


def _fmt(v):
    """Format a number for display."""
    if abs(v) >= 1_000_000:
        return f"${v/1_000_000:,.1f}M"
    if abs(v) >= 10_000:
        return f"${v/1_000:,.1f}K"
    if abs(v) >= 100:
        return f"{v:,.0f}"
    if 0 < abs(v) < 1:
        return f"{v*100:.1f}%"
    return f"{v:,.2f}"


def _num_stats(df):
    """Return list of dicts with stats per numeric column."""
    out = []
    for col in df.select_dtypes(include="number").columns:
        out.append({
            "name": col,
            "min": df[col].min(),
            "max": df[col].max(),
            "mean": df[col].mean(),
            "median": df[col].median(),
            "nulls": int(df[col].isna().sum()),
        })
    return out


class ExportEngine:

    # ================================================================ #
    #  PowerPoint                                                       #
    # ================================================================ #
    def generate_pptx(self, df, title, output_path=None):
        if output_path is None:
            output_path = os.path.join(tempfile.gettempdir(), f"{title}.pptx")
        try:
            print(f"[ExportEngine] Generating PowerPoint: {title}")
            insights = _compute_insights(df)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            bg_rgb = RGBColor(0x0D, 0x0D, 0x0D)
            surface_rgb = RGBColor(0x1A, 0x1A, 0x1A)
            yellow_rgb = RGBColor(0xFF, 0xE6, 0x00)
            white_rgb = RGBColor(0xFF, 0xFF, 0xFF)
            text2_rgb = RGBColor(0xB0, 0xB0, 0xB0)
            muted_rgb = RGBColor(0x80, 0x80, 0x80)
            border_rgb = RGBColor(0x33, 0x33, 0x33)

            def _set_bg(slide):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = bg_rgb

            def _add_text(slide, left, top, width, height, text,
                          size=18, color=white_rgb, bold=False,
                          align=PP_ALIGN.LEFT, line_spacing=1.2):
                txBox = slide.shapes.add_textbox(
                    Inches(left), Inches(top), Inches(width), Inches(height)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(text)
                p.font.size = Pt(size)
                p.font.color.rgb = color
                p.font.bold = bold
                p.alignment = align
                p.space_after = Pt(4)

            def _add_accent_bar(slide):
                """Yellow accent bar at top of slide."""
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    Inches(0), Inches(0),
                    Inches(13.333), Inches(0.06),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = yellow_rgb
                shape.line.fill.background()

            def _add_wu_badge(slide, top=0.3):
                """WESTERN UNION badge in top-right."""
                txBox = slide.shapes.add_textbox(
                    Inches(10.5), Inches(top), Inches(2.5), Inches(0.4)
                )
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = "WESTERN UNION"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                # Yellow background via shape behind
                badge = slide.shapes.add_shape(
                    1, Inches(10.5), Inches(top),
                    Inches(2.3), Inches(0.35),
                )
                badge.fill.solid()
                badge.fill.fore_color.rgb = yellow_rgb
                badge.line.fill.background()
                # Move badge behind text
                slide.shapes._spTree.remove(badge._element)
                slide.shapes._spTree.insert(2, badge._element)

            # ── Slide 1: Title ──
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(sl)
            _add_accent_bar(sl)
            _add_wu_badge(sl)
            _add_text(sl, 1, 2.0, 11, 1.5, title,
                      size=40, color=white_rgb, bold=True, align=PP_ALIGN.CENTER)
            _add_text(sl, 1, 3.8, 11, 0.5,
                      f"{insights['rows']:,} records  |  {insights['cols']} fields  |  "
                      f"{insights['numeric_count']} measures  |  {insights['categorical_count']} dimensions",
                      size=16, color=text2_rgb, align=PP_ALIGN.CENTER)
            _add_text(sl, 1, 4.8, 11, 0.5,
                      datetime.now().strftime("%B %d, %Y"),
                      size=14, color=muted_rgb, align=PP_ALIGN.CENTER)
            _add_text(sl, 1, 5.5, 11, 0.5,
                      "Built by Dr. Data -- Western Union Analytics",
                      size=12, color=muted_rgb, align=PP_ALIGN.CENTER)

            # ── Slide 2: Key Metrics ──
            if insights["kpis"]:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                _set_bg(sl)
                _add_accent_bar(sl)
                _add_wu_badge(sl)
                _add_text(sl, 0.8, 0.4, 8, 0.6, "Key Metrics",
                          size=28, color=yellow_rgb, bold=True)

                for i, kpi in enumerate(insights["kpis"][:4]):
                    col_x = 0.8 + i * 3.1
                    # KPI box background
                    box = sl.shapes.add_shape(
                        1, Inches(col_x), Inches(1.5),
                        Inches(2.8), Inches(2.2),
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = surface_rgb
                    box.line.color.rgb = border_rgb
                    box.line.width = Pt(1)

                    _add_text(sl, col_x + 0.2, 1.7, 2.4, 0.8,
                              kpi["fmt_total"],
                              size=32, color=yellow_rgb, bold=True,
                              align=PP_ALIGN.CENTER)
                    _add_text(sl, col_x + 0.2, 2.5, 2.4, 0.3,
                              kpi["label"],
                              size=12, color=text2_rgb, align=PP_ALIGN.CENTER)
                    _add_text(sl, col_x + 0.2, 2.9, 2.4, 0.3,
                              f"Avg: {kpi['fmt_avg']}",
                              size=11, color=muted_rgb, align=PP_ALIGN.CENTER)

                # Findings below KPIs
                if insights["findings"]:
                    finding_text = "\n".join(
                        f"  {f}" for f in insights["findings"][:4]
                    )
                    _add_text(sl, 0.8, 4.2, 11, 2.5, finding_text,
                              size=16, color=text2_rgb)

            # ── Slide 3: Top Performers ──
            if insights["top_performers"]:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                _set_bg(sl)
                _add_accent_bar(sl)
                _add_wu_badge(sl)
                _add_text(sl, 0.8, 0.4, 8, 0.6, "Top Performers",
                          size=28, color=yellow_rgb, bold=True)

                y_pos = 1.5
                for tp in insights["top_performers"][:3]:
                    _add_text(sl, 0.8, y_pos, 11, 0.4,
                              f"{tp['metric']} by {tp['dimension']}",
                              size=18, color=white_rgb, bold=True)
                    y_pos += 0.5
                    _add_text(sl, 1.2, y_pos, 10, 0.4,
                              f"#1: {tp['top']}  --  {tp['top_val']} ({tp['top_pct']}% of total)",
                              size=16, color=yellow_rgb)
                    y_pos += 0.45
                    if tp["top3"]:
                        _add_text(sl, 1.2, y_pos, 10, 0.4,
                                  f"Top 3: {', '.join(tp['top3'])} = {tp['top3_pct']}% combined  "
                                  f"({tp['count']} total {tp['dimension'].lower()}s)",
                                  size=14, color=text2_rgb)
                        y_pos += 0.6

            # ── Slide 4: Composition ──
            if insights["composition"]:
                sl = prs.slides.add_slide(prs.slide_layouts[6])
                _set_bg(sl)
                _add_accent_bar(sl)
                _add_wu_badge(sl)
                _add_text(sl, 0.8, 0.4, 8, 0.6, "Composition Breakdown",
                          size=28, color=yellow_rgb, bold=True)

                y_pos = 1.4
                for comp in insights["composition"][:2]:
                    _add_text(sl, 0.8, y_pos, 11, 0.4,
                              f"{comp['metric']} by {comp['dimension']}",
                              size=18, color=white_rgb, bold=True)
                    y_pos += 0.55
                    for item in comp["breakdown"][:8]:
                        # Visual bar
                        bar_width = max(0.3, item["pct"] / 100 * 8)
                        bar = sl.shapes.add_shape(
                            1, Inches(3.0), Inches(y_pos + 0.05),
                            Inches(bar_width), Inches(0.25),
                        )
                        bar.fill.solid()
                        bar.fill.fore_color.rgb = yellow_rgb
                        bar.line.fill.background()

                        _add_text(sl, 0.8, y_pos, 2.0, 0.35,
                                  item["name"], size=13, color=white_rgb)
                        _add_text(sl, 3.0 + bar_width + 0.2, y_pos, 3, 0.35,
                                  f"{item['val']}  ({item['pct']}%)",
                                  size=12, color=text2_rgb)
                        y_pos += 0.38

            # ── Slide 5: Data Quality ──
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(sl)
            _add_accent_bar(sl)
            _add_wu_badge(sl)
            _add_text(sl, 0.8, 0.4, 8, 0.6, "Data Quality Assessment",
                      size=28, color=yellow_rgb, bold=True)

            quality_text = (
                f"Records: {insights['rows']:,}\n"
                f"Fields: {insights['cols']}\n"
                f"Overall null rate: {insights['null_pct']}%\n\n"
            )
            quality_text += "\n".join(
                f"  {note}" for note in insights["quality_notes"]
            )
            _add_text(sl, 0.8, 1.4, 11, 4, quality_text,
                      size=18, color=text2_rgb)

            # Stats table as appendix info
            stats = _num_stats(df)
            if stats:
                stat_lines = []
                for s in stats[:8]:
                    stat_lines.append(
                        f"{s['name'][:20]:20s}  "
                        f"Min: {s['min']:>12,.2f}  Max: {s['max']:>12,.2f}  "
                        f"Avg: {s['mean']:>12,.2f}"
                    )
                _add_text(sl, 0.8, 3.8, 12, 3, "\n".join(stat_lines),
                          size=11, color=muted_rgb)

            # ── Slide 6: Closing ──
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(sl)
            _add_accent_bar(sl)
            _add_wu_badge(sl, top=0.3)
            _add_text(sl, 1, 2.5, 11, 1, "Thank You",
                      size=36, color=white_rgb, bold=True, align=PP_ALIGN.CENTER)
            _add_text(sl, 1, 3.8, 11, 0.5,
                      "Built by Dr. Data -- Western Union Analytics",
                      size=16, color=muted_rgb, align=PP_ALIGN.CENTER)
            _add_text(sl, 1, 4.5, 11, 0.5,
                      datetime.now().strftime("%B %d, %Y"),
                      size=14, color=muted_rgb, align=PP_ALIGN.CENTER)

            prs.save(output_path)
            print(f"[ExportEngine] PowerPoint saved: {output_path}")
            return output_path
        except Exception as e:
            print(f"[ExportEngine] PowerPoint failed: {e}")
            import traceback; traceback.print_exc()
            return None

    # ================================================================ #
    #  Word Document                                                    #
    # ================================================================ #
    def generate_docx(self, df, title, output_path=None):
        if output_path is None:
            output_path = os.path.join(tempfile.gettempdir(), f"{title}.docx")
        try:
            print(f"[ExportEngine] Generating Word doc: {title}")
            insights = _compute_insights(df)
            doc = Document()

            # Style defaults
            style = doc.styles["Normal"]
            font = style.font
            font.name = "Calibri"
            font.size = DocxPt(11)

            # ── Title ──
            h = doc.add_heading(title, level=0)
            for run in h.runs:
                run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)
                run.font.size = DocxPt(28)

            p = doc.add_paragraph()
            run = p.add_run("WESTERN UNION")
            run.bold = True
            run.font.size = DocxPt(10)
            run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

            doc.add_paragraph(
                f"{datetime.now().strftime('%B %d, %Y')}  |  "
                f"{insights['rows']:,} records  |  {insights['cols']} fields"
            )
            doc.add_paragraph("")

            # ── Executive Summary ──
            h2 = doc.add_heading("Executive Summary", level=1)
            for run in h2.runs:
                run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

            if insights["findings"]:
                for finding in insights["findings"]:
                    p = doc.add_paragraph(style="List Bullet")
                    p.text = finding
            else:
                doc.add_paragraph(
                    f"Analysis of {insights['rows']:,} records across "
                    f"{insights['cols']} fields. "
                    f"{insights['numeric_count']} numeric measures and "
                    f"{insights['categorical_count']} categorical dimensions identified."
                )

            doc.add_paragraph("")

            # ── Key Metrics ──
            if insights["kpis"]:
                h2 = doc.add_heading("Key Metrics", level=1)
                for run in h2.runs:
                    run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

                table = doc.add_table(rows=1, cols=5)
                table.style = "Table Grid"
                hdr = table.rows[0].cells
                for i, label in enumerate(
                    ["Metric", "Total", "Average", "Min", "Max"]
                ):
                    hdr[i].text = label
                    for p in hdr[i].paragraphs:
                        p.runs[0].bold = True if p.runs else False

                for kpi in insights["kpis"]:
                    row = table.add_row().cells
                    row[0].text = kpi["label"]
                    row[1].text = kpi["fmt_total"]
                    row[2].text = kpi["fmt_avg"]
                    row[3].text = _fmt(kpi["min"])
                    row[4].text = _fmt(kpi["max"])

                doc.add_paragraph("")

            # ── Top Performers ──
            if insights["top_performers"]:
                h2 = doc.add_heading("Top Performers", level=1)
                for run in h2.runs:
                    run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

                for tp in insights["top_performers"]:
                    doc.add_heading(
                        f"{tp['metric']} by {tp['dimension']}", level=2
                    )
                    p = doc.add_paragraph(style="List Bullet")
                    p.text = (
                        f"#1: {tp['top']} at {tp['top_val']} "
                        f"({tp['top_pct']}% of total)"
                    )
                    if tp["top3"]:
                        p = doc.add_paragraph(style="List Bullet")
                        p.text = (
                            f"Top 3: {', '.join(tp['top3'])} = "
                            f"{tp['top3_pct']}% combined"
                        )
                    p = doc.add_paragraph(style="List Bullet")
                    p.text = f"Total {tp['dimension'].lower()}s: {tp['count']}"

                doc.add_paragraph("")

            # ── Composition Breakdown ──
            if insights["composition"]:
                h2 = doc.add_heading("Composition Breakdown", level=1)
                for run in h2.runs:
                    run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

                for comp in insights["composition"]:
                    doc.add_heading(
                        f"{comp['metric']} by {comp['dimension']}", level=2
                    )
                    table = doc.add_table(rows=1, cols=3)
                    table.style = "Table Grid"
                    hdr = table.rows[0].cells
                    hdr[0].text = comp["dimension"]
                    hdr[1].text = comp["metric"]
                    hdr[2].text = "Share"
                    for item in comp["breakdown"][:10]:
                        row = table.add_row().cells
                        row[0].text = item["name"]
                        row[1].text = item["val"]
                        row[2].text = f"{item['pct']}%"

                doc.add_paragraph("")

            # ── Data Quality ──
            h2 = doc.add_heading("Data Quality Assessment", level=1)
            for run in h2.runs:
                run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

            p = doc.add_paragraph()
            p.text = (
                f"Dataset: {insights['rows']:,} rows x {insights['cols']} columns. "
                f"Overall null rate: {insights['null_pct']}%."
            )
            for note in insights["quality_notes"]:
                p = doc.add_paragraph(style="List Bullet")
                p.text = note

            doc.add_paragraph("")

            # ── Statistical Reference ──
            stats = _num_stats(df)
            if stats:
                h2 = doc.add_heading("Statistical Reference", level=1)
                for run in h2.runs:
                    run.font.color.rgb = DocxRGB(0xFF, 0xE6, 0x00)

                table = doc.add_table(rows=1, cols=6)
                table.style = "Table Grid"
                hdr = table.rows[0].cells
                for i, label in enumerate(
                    ["Column", "Min", "Max", "Mean", "Median", "Nulls"]
                ):
                    hdr[i].text = label
                for s in stats[:15]:
                    row = table.add_row().cells
                    row[0].text = s["name"][:25]
                    row[1].text = f"{s['min']:,.2f}"
                    row[2].text = f"{s['max']:,.2f}"
                    row[3].text = f"{s['mean']:,.2f}"
                    row[4].text = f"{s['median']:,.2f}"
                    row[5].text = str(s["nulls"])

            # ── Footer ──
            doc.add_paragraph("")
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run("Built by Dr. Data -- Western Union Analytics")
            run.font.size = DocxPt(9)
            run.font.color.rgb = DocxRGB(0x80, 0x80, 0x80)

            doc.save(output_path)
            print(f"[ExportEngine] Word doc saved: {output_path}")
            return output_path
        except Exception as e:
            print(f"[ExportEngine] Word doc failed: {e}")
            import traceback; traceback.print_exc()
            return None

    # ================================================================ #
    #  PDF                                                              #
    # ================================================================ #
    def generate_pdf(self, df, title, output_path=None):
        if output_path is None:
            output_path = os.path.join(tempfile.gettempdir(), f"{title}.pdf")
        try:
            print(f"[ExportEngine] Generating PDF: {title}")
            insights = _compute_insights(df)

            doc = SimpleDocTemplate(
                output_path, pagesize=letter,
                topMargin=0.5 * inch, bottomMargin=0.5 * inch,
                leftMargin=0.5 * inch, rightMargin=0.5 * inch,
            )
            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                "DrTitle", parent=styles["Title"],
                fontSize=22, textColor=HexColor(WU_YELLOW), spaceAfter=4,
            )
            subtitle_style = ParagraphStyle(
                "DrSubtitle", parent=styles["Normal"],
                fontSize=10, textColor=HexColor(WU_MUTED), spaceAfter=12,
            )
            body_style = ParagraphStyle(
                "DrBody", parent=styles["Normal"],
                fontSize=10, textColor=HexColor("#333333"), spaceAfter=8,
                leading=14,
            )
            heading_style = ParagraphStyle(
                "DrHeading", parent=styles["Heading2"],
                fontSize=14, textColor=HexColor(WU_YELLOW), spaceAfter=6,
                backColor=HexColor(WU_BLACK),
            )
            bullet_style = ParagraphStyle(
                "DrBullet", parent=styles["Normal"],
                fontSize=10, textColor=HexColor("#333333"),
                spaceAfter=4, leftIndent=20, bulletIndent=10,
                leading=14,
            )

            elements = []
            elements.append(Paragraph(title, title_style))
            elements.append(Paragraph(
                f"WESTERN UNION  |  {datetime.now().strftime('%B %d, %Y')}  |  "
                f"{insights['rows']:,} records  |  {insights['cols']} fields",
                subtitle_style,
            ))
            elements.append(Spacer(1, 12))

            # Executive summary
            elements.append(Paragraph("Executive Summary", heading_style))
            for finding in insights["findings"][:5]:
                elements.append(Paragraph(f"\u2022  {finding}", bullet_style))
            if not insights["findings"]:
                elements.append(Paragraph(
                    f"Analysis of {insights['rows']:,} records across "
                    f"{insights['cols']} fields.",
                    body_style,
                ))
            elements.append(Spacer(1, 12))

            # Key metrics table
            if insights["kpis"]:
                elements.append(Paragraph("Key Metrics", heading_style))
                header = ["Metric", "Total", "Average", "Min", "Max"]
                data = [header]
                for kpi in insights["kpis"]:
                    data.append([
                        kpi["label"], kpi["fmt_total"], kpi["fmt_avg"],
                        _fmt(kpi["min"]), _fmt(kpi["max"]),
                    ])
                t = Table(data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), HexColor(WU_BLACK)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor(WU_YELLOW)),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor(WU_BORDER)),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                     [HexColor("#FFFFFF"), HexColor("#F5F5F5")]),
                    ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 12))

            # Top performers
            if insights["top_performers"]:
                elements.append(Paragraph("Top Performers", heading_style))
                for tp in insights["top_performers"][:3]:
                    elements.append(Paragraph(
                        f"<b>{tp['metric']} by {tp['dimension']}:</b>  "
                        f"#1 is {tp['top']} at {tp['top_val']} ({tp['top_pct']}%). "
                        f"Top 3 = {tp['top3_pct']}% combined.",
                        body_style,
                    ))
                elements.append(Spacer(1, 12))

            # Stats table
            stats = _num_stats(df)
            if stats:
                elements.append(Paragraph("Statistical Reference", heading_style))
                header = ["Column", "Min", "Max", "Mean", "Median", "Nulls"]
                data = [header]
                for s in stats[:15]:
                    data.append([
                        s["name"][:20],
                        f"{s['min']:,.2f}", f"{s['max']:,.2f}",
                        f"{s['mean']:,.2f}", f"{s['median']:,.2f}",
                        str(s["nulls"]),
                    ])
                t = Table(data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), HexColor(WU_BLACK)),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor(WU_YELLOW)),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor(WU_BORDER)),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                     [HexColor("#FFFFFF"), HexColor("#F5F5F5")]),
                    ("TEXTCOLOR", (0, 1), (-1, -1), HexColor("#333333")),
                    ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 12))

            # Data quality
            elements.append(Paragraph("Data Quality", heading_style))
            elements.append(Paragraph(
                f"Overall null rate: {insights['null_pct']}%", body_style
            ))
            for note in insights["quality_notes"]:
                elements.append(Paragraph(f"\u2022  {note}", bullet_style))

            # Footer
            elements.append(Spacer(1, 20))
            footer_style = ParagraphStyle(
                "DrFooter", parent=styles["Normal"],
                fontSize=8, textColor=HexColor(WU_MUTED),
                alignment=1,  # center
            )
            elements.append(Paragraph(
                "Built by Dr. Data -- Western Union Analytics", footer_style
            ))

            doc.build(elements)
            print(f"[ExportEngine] PDF saved: {output_path}")
            return output_path
        except Exception as e:
            print(f"[ExportEngine] PDF failed: {e}")
            import traceback; traceback.print_exc()
            return None

    # ================================================================ #
    #  All formats                                                      #
    # ================================================================ #
    def generate_all(self, df, title, output_dir=None):
        if output_dir is None:
            output_dir = tempfile.gettempdir()
        os.makedirs(output_dir, exist_ok=True)
        safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in title)
        results = {}
        results["pptx"] = self.generate_pptx(
            df, title, os.path.join(output_dir, f"{safe}.pptx")
        )
        results["pdf"] = self.generate_pdf(
            df, title, os.path.join(output_dir, f"{safe}.pdf")
        )
        results["docx"] = self.generate_docx(
            df, title, os.path.join(output_dir, f"{safe}.docx")
        )
        return results
